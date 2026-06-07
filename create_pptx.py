from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(10, 22, 40)
MID_BG = RGBColor(13, 35, 71)
CYAN = RGBColor(0, 220, 180)
BLUE = RGBColor(0, 160, 255)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(176, 196, 216)
ORANGE = RGBColor(255, 165, 0)
GREEN = RGBColor(52, 168, 83)
RED = RGBColor(234, 67, 53)
PURPLE = RGBColor(155, 89, 182)
YELLOW = RGBColor(251, 188, 4)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color=MID_BG, left=0, top=0, width=None, height=None):
    if width is None: width = prs.slide_width
    if height is None: height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16, color=WHITE, bold_first=True, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def add_card(slide, left, top, width, height, title, content_lines, card_color=CYAN, title_color=DARK_BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = card_color
    shape.line.width = Pt(2)
    
    # Title bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = card_color
    bar.line.fill.background()
    
    add_textbox(slide, left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), Inches(0.45),
                title, 14, title_color, True, PP_ALIGN.CENTER)
    
    add_multiline_textbox(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), height - Inches(0.65),
                          content_lines, 12, DARK_BG, False)
    return shape

def add_accent_line(slide, left, top, width, color=CYAN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BG)

# Decorative shapes
add_shape_bg(slide, MID_BG, Inches(0), Inches(0), Inches(13.333), Inches(4.5))
add_accent_line(slide, Inches(1), Inches(4.5), Inches(11.333), CYAN)

# Dots pattern
for x in range(20, 800, 60):
    for y in range(20, 180, 60):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x/100), Inches(y/100), Pt(2), Pt(2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = RGBColor(255, 255, 255)
        dot.fill.fore_color.brightness = 0.9
        dot.line.fill.background()

add_textbox(slide, Inches(1), Inches(1.2), Inches(11.333), Inches(1),
            'INDEXATION ET RÉFÉRENCEMENT WEB', 40, WHITE, True, PP_ALIGN.LEFT)

add_textbox(slide, Inches(1), Inches(2.3), Inches(11.333), Inches(0.6),
            'Guide Complet et Approfondi — SEO Search Engine Optimization', 20, CYAN, False, PP_ALIGN.LEFT)

# Tags
tags = ['On-Page SEO', 'Techniques', 'Outils', 'Stratégie']
for i, tag in enumerate(tags):
    x = Inches(1 + i * 2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.3), Inches(2), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.fill.fore_color.brightness = 0.85
    shape.line.fill.background()
    add_textbox(slide, x, Inches(3.33), Inches(2), Inches(0.4), tag, 13, DARK_BG, True, PP_ALIGN.CENTER)

# Institution info
add_textbox(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(0.5),
            'Institut Supérieur du Numérique — SupNum, Nouakchott', 16, LIGHT_GRAY, False)
add_textbox(slide, Inches(1), Inches(5.7), Inches(11.333), Inches(0.5),
            'Spécialité : Développement Web et Multimédia', 14, LIGHT_GRAY, False)

add_textbox(slide, Inches(1), Inches(6.4), Inches(11.333), Inches(0.5),
            'Version 1.0 — Juin 2026', 12, LIGHT_GRAY, False, PP_ALIGN.RIGHT)


# ===== SLIDE 2: Plan =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), CYAN)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8), 'Plan de la présentation', 32, WHITE, True)

parts = [
    ('Partie I', 'Fondamentaux des moteurs de recherche', 'Histoire, fonctionnement, IA', '3 chapitres'),
    ('Partie II', 'L\'Indexation', 'Crawling, robots.txt, sitemaps, budget de crawl', '4 chapitres'),
    ('Partie III', 'SEO On-Page', 'Mots-clés, intention, balises HTML, images, maillage', '5 chapitres'),
    ('Partie IV', 'SEO Technique Avancé', 'Core Web Vitals, performance, mobile, schema, JS', '6 chapitres'),
    ('Partie V', 'SEO Off-Page', 'Netlinking, link building, E-E-A-T', '3 chapitres'),
    ('Partie VI', 'SEO Spécialisé', 'Local, e-commerce, CMS, vidéo, vocal, A/B testing', '6 chapitres'),
    ('Partie VII', 'Mesure et Analyse', 'GSC, GA4, KPIs et reporting', '3 chapitres'),
    ('Partie VIII', 'Stratégie et Tendances', 'Audit, plan d\'action, IA, futur du SEO', '4 chapitres'),
]

for i, (num, title, desc, count) in enumerate(parts):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.5 + row * 1.4)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.color.rgb = CYAN if i % 2 == 0 else BLUE
    shape.line.width = Pt(1.5)
    
    add_textbox(slide, left + Inches(0.2), top + Inches(0.05), Inches(1.2), Inches(0.35),
                num, 11, CYAN, True)
    add_textbox(slide, left + Inches(1.3), top + Inches(0.05), Inches(4.3), Inches(0.35),
                title, 15, WHITE, True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.45), Inches(4.5), Inches(0.35),
                desc, 11, LIGHT_GRAY, False)
    add_textbox(slide, left + Inches(4.5), top + Inches(0.85), Inches(1.2), Inches(0.3),
                count, 10, CYAN, False, PP_ALIGN.RIGHT)


# ===== SLIDE 3: Méthodologie =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), GREEN)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8), 'Méthodologie de rédaction', 32, WHITE, True)

items = [
    ('🎯', 'Approche pédagogique', 'Structure homogène par chapitre : objectifs → théorie → illustrations → exercices → synthèse'),
    ('📚', 'Sources officielles', 'Documentation Google Search Central, ingénieurs Google (John Mueller, Gary Illyes), recherches académiques'),
    ('🎓', 'Public cible', 'Étudiants Développement Web et Multimédia (Licence/Master), professionnels du numérique'),
    ('🔧', 'Prérequis', 'Connaissances de base en HTML/CSS, familiarité avec le Web, notions de marketing digital'),
    ('📊', 'Volume horaire', '56 heures (30h cours + 26h TD/TP), réparties sur 14 semaines'),
]

for i, (icon, title, desc) in enumerate(items):
    top = Inches(1.5 + i * 1.1)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), top, Inches(11.333), Inches(0.95))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.fill.background()
    
    add_textbox(slide, Inches(1.2), top + Inches(0.05), Inches(0.5), Inches(0.85),
                icon, 24, WHITE, False, PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.8), top + Inches(0.05), Inches(3), Inches(0.4),
                title, 16, GREEN, True)
    add_textbox(slide, Inches(1.8), top + Inches(0.45), Inches(10), Inches(0.45),
                desc, 12, LIGHT_GRAY, False)


# ===== SLIDE 4: Part I - Fondamentaux =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), CYAN)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Partie I — Fondamentaux des moteurs de recherche', 28, WHITE, True)

cards = [
    ('Histoire & évolution', ['1990 : Archie (1er moteur)', '1998 : Google et le PageRank', 'Index inversé : révolution technique', 'Mises à jour algorithmiques majeures', '91% de parts de marché pour Google'], CYAN),
    ('Fonctionnement technique', ['3 étapes : Crawling → Indexation → Ranking', 'Googlebot explore le web via des liens', 'L\'index inversé associe mot-clé → documents', 'Le PageRank mesure l\'autorité par les liens', 'Centaines de critères de classement'], BLUE),
    ('IA dans les moteurs', ['RankBrain (2015) : apprentissage automatique', 'BERT (2019) : compréhension du langage', 'MUM (2021) : modèle multimodal', 'Compréhension sémantique des requêtes', 'Traitement du langage naturel avancé'], PURPLE),
]

for i, (title, items, color) in enumerate(cards):
    left = Inches(0.8 + i * 4.1)
    add_card(slide, left, Inches(1.5), Inches(3.8), Inches(5.5), title, items, color)


# ===== SLIDE 5: Part II - Indexation =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), ORANGE)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Partie II — L\'Indexation, pilier technique', 28, WHITE, True)

topics = [
    ('Pipeline d\'indexation', ['Analyse → Extraction → Stockage → Mise à jour', 'Vérification via Google Search Console', 'Problèmes : erreurs de crawl, soft 404, pages orphelines'], ORANGE),
    ('Contrôle du crawling', ['robots.txt : directives pour Googlebot', 'Sitemap XML : guide vers les pages importantes', 'Balises meta robots : noindex, nofollow'], ORANGE),
    ('Budget de crawl', ['Pages explorées par période', 'Facteurs : qualité, popularité, taille du site', 'Optimisation : supprimer le contenu faible valeur'], ORANGE),
]

for i, (title, items, color) in enumerate(topics):
    left = Inches(0.8 + i * 4.1)
    add_card(slide, left, Inches(1.5), Inches(3.8), Inches(4.5), title, items, color)

add_textbox(slide, Inches(1), Inches(6.3), Inches(11.333), Inches(0.6),
            '💡 Conseil : Utilisez le rapport de couverture GSC pour diagnostiquer vos problèmes d\'indexation',
            14, YELLOW, False)


# ===== SLIDE 6: Part III - On-Page =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), GREEN)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Partie III — SEO On-Page', 28, WHITE, True)

onpage = [
    ('Recherche de mots-clés', 'Fondation de toute stratégie SEO\nTête de réseau, corps, longue traîne\nOutils : SEMrush, Ahrefs, KW Finder\nVolume, difficulté, CPC comme métriques'),
    ('Intention de recherche', '4 types : Informationnelle, Navigationnelle\nCommerciale, Transactionnelle\nLe contenu doit correspondre à l\'intention\nCartographie du parcours utilisateur'),
    ('Balises HTML', 'Title tag : 2e facteur On-Page le + important\nMeta description : influence le CTR\nHiérarchie H1-H6 : structure du contenu'),
    ('Images & maillage', 'Alt text : accessibilité + SEO\nFormats modernes : WebP, AVIF\nTopical clusters : pilier + articles satellites\nMaillage interne : distribue l\'autorité'),
]

for i, (title, desc) in enumerate(onpage):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(1.5 + row * 2.8)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.5)
    
    add_textbox(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.3), Inches(0.4),
                title, 17, GREEN, True)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.55), Inches(5.3), Inches(1.8),
                desc, 13, LIGHT_GRAY, False)


# ===== SLIDE 7: Part IV - Technique Avancé =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), RED)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Partie IV — SEO Technique Avancé', 28, WHITE, True)

# Core Web Vitals highlight
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = MID_BG
shape.line.color.rgb = RED
shape.line.width = Pt(2)

cwv = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(1.5), Inches(1.2))
cwv.fill.solid()
cwv.fill.fore_color.rgb = RED
cwv.line.fill.background()
add_textbox(slide, Inches(0.9), Inches(1.7), Inches(1.3), Inches(0.8), 'Core\nWeb\nVitals', 18, WHITE, True, PP_ALIGN.CENTER)

add_textbox(slide, Inches(2.5), Inches(1.6), Inches(3), Inches(0.4),
            'LCP  < 2,5s', 18, RGBColor(52, 168, 83), True)
add_textbox(slide, Inches(2.5), Inches(2.0), Inches(3), Inches(0.3),
            'Largest Contentful Paint', 12, LIGHT_GRAY, False)

add_textbox(slide, Inches(5.5), Inches(1.6), Inches(3), Inches(0.4),
            'INP  < 200ms', 18, RGBColor(251, 188, 4), True)
add_textbox(slide, Inches(5.5), Inches(2.0), Inches(3), Inches(0.3),
            'Interaction to Next Paint', 12, LIGHT_GRAY, False)

add_textbox(slide, Inches(8.5), Inches(1.6), Inches(3.5), Inches(0.4),
            'CLS  < 0,1', 18, RGBColor(66, 133, 244), True)
add_textbox(slide, Inches(8.5), Inches(2.0), Inches(3.5), Inches(0.3),
            'Cumulative Layout Shift', 12, LIGHT_GRAY, False)

# Other topics
tech_topics = [
    ('Performance web', 'CDN, HTTP/2/3, caching\nOptimisation images (WebP, AVIF)\nCode splitting, tree shaking\nCritical CSS, lazy loading'),
    ('Mobile-First', 'Google utilise la version mobile\nDesign responsive obligatoire\nPerformances mobiles essentielles'),
    ('Schema Markup', 'Schema.org, JSON-LD\nRich snippets\nTypes : Article, FAQ, Product'),
    ('JavaScript SEO', 'Défis : rendu JS, indexation\nSSR, SSG, hydration\nTest dans Google Search Console'),
]

for i, (title, desc) in enumerate(tech_topics):
    row = i // 2
    col = i % 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(3.0 + row * 2.0)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.8), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.color.rgb = RED
    shape.line.width = Pt(1)
    
    add_textbox(slide, left + Inches(0.3), top + Inches(0.1), Inches(5.3), Inches(0.4),
                title, 16, RED, True)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.5), Inches(5.3), Inches(1.2),
                desc, 12, LIGHT_GRAY, False)


# ===== SLIDE 8: Part V - Off-Page =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), PURPLE)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Partie V — SEO Off-Page et Autorité', 28, WHITE, True)

offpage = [
    ('Netlinking & Backlinks', ['Facteur de classement majeur', 'Dofollow : transmet l\'autorité', 'Nofollow : ne transmet pas', 'Qualité > Quantité', 'Profil de liens naturel'], PURPLE),
    ('Link Building avancé', ['Digital PR : relations presse', 'Linkable assets : contenu magnétique', 'Broken link building', 'Guest blogging ciblé', 'Brand mentions'], PURPLE),
    ('E-E-A-T', ['Expérience, Expertise, Autorité, Confiance', 'Cadre d\'évaluation qualité Google', 'YMYL : sites à impact fort', 'Auteurs identifiés et qualifiés', 'Transparence et fiabilité'], PURPLE),
]

for i, (title, items, color) in enumerate(offpage):
    left = Inches(0.8 + i * 4.1)
    add_card(slide, left, Inches(1.5), Inches(3.8), Inches(5), title, items, color)


# ===== SLIDE 9: Part VI - Spécialisé =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), ORANGE)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Partie VI — SEO Spécialisé', 28, WHITE, True)

spec = [
    ('SEO Local', 'Google Business Profile\nCitations NAP, annuaires\nAvis clients'),
    ('Migration SEO', '301 vs 302, mapping URL\nChecklist pré-migration\nSuivi post-migration (3 mois)'),
    ('E-commerce', 'Architecture optimale\nContenu dupliqué\nFiches produits uniques'),
    ('CMS populaires', 'WordPress, Shopify, Webflow\nPlugins SEO (Rank Math, Yoast)\nPerformance par CMS'),
    ('Vidéo & Vocal', 'YouTube = 2e moteur mondial\nFeatured snippets pour vocal\nDonnées structurées vidéo'),
    ('A/B Testing', 'Split testing SEO\nSignification statistique\nPriorité : titles, schema'),
]

for i, (title, desc) in enumerate(spec):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(1.5 + row * 2.8)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.color.rgb = ORANGE
    shape.line.width = Pt(1.5)
    
    add_textbox(slide, left + Inches(0.2), top + Inches(0.15), Inches(3.4), Inches(0.4),
                title, 15, ORANGE, True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.55), Inches(3.4), Inches(1.8),
                desc, 12, LIGHT_GRAY, False)


# ===== SLIDE 10: Part VII-VIII - Mesure & Stratégie =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), CYAN)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Parties VII & VIII — Mesure, Stratégie et Tendances', 26, WHITE, True)

# Left column: Mesure
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5),
            '📊 Mesure et Analyse', 20, GREEN, True)

measure_items = [
    'Google Search Console : indexation, performances, erreurs',
    'Google Analytics 4 : trafic, comportement, conversions',
    'KPIs SEO : trafic organique, positions, CTR, taux d\'indexation',
    'Tableau de bord SEO adapté aux objectifs business',
    'Reporting régulier pour ajuster la stratégie',
]
add_multiline_textbox(slide, Inches(0.8), Inches(2.1), Inches(5.8), Inches(3), measure_items, 13, LIGHT_GRAY, False)

# Right column: Stratégie
add_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.8), Inches(0.5),
            '🎯 Stratégie et Tendances', 20, BLUE, True)

strat_items = [
    'Audit SEO complet : technique, On-Page, Off-Page',
    'Plan d\'action priorisé par impact et effort',
    'SEO et IA : Google SGE, Search Generative Experience',
    'Tendances : zero-click, Search Everywhere Optimization',
    'Prévisions 2026-2030 : Web3, SEO décentralisé',
    'L\'adaptation aux nouvelles technologies = clé de la pérennité',
]
add_multiline_textbox(slide, Inches(7.2), Inches(2.1), Inches(5.8), Inches(3.5), strat_items, 13, LIGHT_GRAY, False)

# Separator line
add_accent_line(slide, Inches(6.7), Inches(1.5), Pt(3), LIGHT_GRAY)

# Case studies highlight
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.733), Inches(0.5),
            '📋 4 études de cas incluses : Migration site • E-commerce CWV • Stratégie contenu • Expansion internationale',
            15, YELLOW, False, PP_ALIGN.CENTER)


# ===== SLIDE 11: Structure pédagogique =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), GREEN)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Structure pédagogique du manuel', 28, WHITE, True)

ped_items = [
    ('🎯', '38 objectifs pédagogiques', 'Un par chapitre, 3 objectifs mesurables chacun'),
    ('📌', '37 synthèses', 'Points à retenir en fin de chaque chapitre'),
    ('✏️', '4 exercices pratiques', 'Indexation, mots-clés, Core Web Vitals, audit complet'),
    ('📖', 'Glossaire complet', 'Définitions de tous les termes techniques SEO'),
    ('📋', '3 annexes', 'Checklist technique, guide des outils, ressources'),
    ('📚', 'Bibliographie', 'Sources officielles, blogs, livres de référence'),
]

for i, (icon, title, desc) in enumerate(ped_items):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.1)
    top = Inches(1.5 + row * 2.5)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.5)
    
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), Inches(3.4), Inches(0.6),
                f'{icon}  {title}', 16, GREEN, True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.7), Inches(3.4), Inches(1.3),
                desc, 12, LIGHT_GRAY, False)


# ===== SLIDE 12: Appendices =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), PURPLE)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Annexes et ressources', 28, WHITE, True)

app_items = [
    ('Annexe A — Checklist technique SEO', [
        '7 catégories : crawling, architecture, performance, mobile, schema, sécurité, international',
        '30+ points de vérification pour auditer un site',
    ]),
    ('Annexe B — Guide des outils SEO', [
        '15 outils avec prix : GSC, GA4, Screaming Frog, Ahrefs, SEMrush, Majestic, Surfer SEO...',
        'Métriques clés à suivre : trafic, positions, CTR, indexation, backlinks',
    ]),
    ('Annexe C — Ressources complémentaires', [
        'Outils de veille : Google Search Status, SEOFrog',
        'Communautés : WebmasterWorld, Reddit r/SEO',
        'Certifications : Google, Ahrefs Academy, SEMrush Academy',
    ]),
]

for i, (title, items) in enumerate(app_items):
    top = Inches(1.5 + i * 1.8)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.color.rgb = PURPLE
    shape.line.width = Pt(1.5)
    
    add_textbox(slide, Inches(1.2), top + Inches(0.1), Inches(11), Inches(0.4),
                title, 16, PURPLE, True)
    for j, item in enumerate(items):
        add_textbox(slide, Inches(1.2), top + Inches(0.5 + j * 0.4), Inches(11), Inches(0.4),
                    f'• {item}', 12, LIGHT_GRAY, False)


# ===== SLIDE 13: Chiffres clés =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_accent_line(slide, Inches(1), Inches(0.2), Inches(3), YELLOW)
add_textbox(slide, Inches(1), Inches(0.4), Inches(11.333), Inches(0.8),
            'Chiffres clés du manuel', 28, WHITE, True)

stats = [
    ('168', 'pages'),
    ('38', 'chapitres'),
    ('8', 'parties'),
    ('37', 'synthèses'),
    ('38', 'objectifs'),
    ('7+', 'schémas SVG'),
    ('15+', 'outils référencés'),
    ('4', 'études de cas'),
]

for i, (num, label) in enumerate(stats):
    row = i // 4
    col = i % 4
    left = Inches(0.8 + col * 3.1)
    top = Inches(1.8 + row * 2.5)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BG
    shape.line.color.rgb = YELLOW
    shape.line.width = Pt(2)
    
    add_textbox(slide, left, top + Inches(0.2), Inches(2.8), Inches(0.9),
                num, 44, YELLOW, True, PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(1.2), Inches(2.8), Inches(0.5),
                label, 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ===== SLIDE 14: Thank you =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BG)
add_shape_bg(slide, MID_BG, Inches(0), Inches(0), Inches(13.333), Inches(4))

add_accent_line(slide, Inches(2), Inches(4.5), Inches(9.333), CYAN)

# decorative circles
for x in range(200, 1200, 80):
    for y in range(30, 160, 60):
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x/100), Inches(y/100), Pt(2), Pt(2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = WHITE
        dot.fill.fore_color.brightness = 0.9
        dot.line.fill.background()

add_textbox(slide, Inches(1), Inches(1.5), Inches(11.333), Inches(1.2),
            'Merci de votre attention', 44, WHITE, True, PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3), Inches(11.333), Inches(0.8),
            'Guide complet disponible sur GitHub', 22, CYAN, False, PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(0.5),
            'Institut Supérieur du Numérique — SupNum, Nouakchott', 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(5.7), Inches(11.333), Inches(0.5),
            'Spécialité : Développement Web et Multimédia — Module : Indexation et Référencement Web', 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.3), Inches(11.333), Inches(0.5),
            'github.com/mohameden19961/rapport-seo', 14, BLUE, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(6.8), Inches(11.333), Inches(0.5),
            'Version 1.0 — Juin 2026', 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# Save
output_path = '/home/abdy/rapport_seo/presentation_seo.pptx'
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Size: {os.path.getsize(output_path) / 1024:.1f} KB")
print(f"Slides: {len(prs.slides)}")
