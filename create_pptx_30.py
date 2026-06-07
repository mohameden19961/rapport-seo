from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(10, 22, 40)
MID_BG = RGBColor(13, 35, 71)
CYAN = RGBColor(0, 220, 180)
BLUE = RGBColor(0, 160, 255)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(176, 196, 216)
ORANGE = RGBColor(255, 165, 0)
GREEN = RGBColor(52, 168, 83)
RED = RGBColor(234, 67, 53)
PURPLE = RGBColor(155, 89, 182)
YELLOW = RGBColor(251, 188, 4)
GOOGLE_BLUE = RGBColor(66, 133, 244)

def bg(slide, color=DARK_BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, alpha=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def rrect(slide, l, t, w, h, color, border=None, bw=Pt(1.5)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border
        s.line.width = bw
    else:
        s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    bx = slide.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    p = bx.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    return bx

def multi(slide, l, t, w, h, lines, sz=13, color=WHITE, bold_first=False, spacing=6):
    bx = slide.shapes.add_textbox(l, t, w, h)
    bx.text_frame.word_wrap = True
    for i, line in enumerate(lines):
        p = bx.text_frame.paragraphs[0] if i == 0 else bx.text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(spacing)
        if bold_first and i == 0:
            p.font.bold = True
    return bx

def accent(slide, l, t, w, color=CYAN):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Pt(3))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def dot_pattern(slide):
    for x in range(20, 800, 60):
        for y in range(20, 180, 60):
            d = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x/100), Inches(y/100), Pt(2), Pt(2))
            d.fill.solid()
            d.fill.fore_color.rgb = WHITE
            d.fill.fore_color.brightness = 0.9
            d.line.fill.background()

def card(slide, l, t, w, h, title, items, color=CYAN):
    rrect(slide, l, t, w, h, WHITE, color)
    bar = rect(slide, l, t, w, Inches(0.45), color)
    txt(slide, l+Inches(0.1), t+Inches(0.03), w-Inches(0.2), Inches(0.4), title, 13, DARK_BG, True, PP_ALIGN.CENTER)
    multi(slide, l+Inches(0.12), t+Inches(0.5), w-Inches(0.24), h-Inches(0.6), items, 11, DARK_BG)

def slide_header(slide, num, title, color=CYAN):
    accent(slide, Inches(1), Inches(0.15), Inches(3), color)
    txt(slide, Inches(1), Inches(0.35), Inches(11.333), Inches(0.7), title, 26, WHITE, True)
    txt(slide, Inches(11.5), Inches(0.35), Inches(1.5), Inches(0.5), f'{num}/30', 12, LIGHT_GRAY, False, PP_ALIGN.RIGHT)

# ===================== SLIDE 1: TITLE =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(4.5), MID_BG)
accent(s, Inches(1), Inches(4.5), Inches(11.333), CYAN)
dot_pattern(s)
txt(s, Inches(1), Inches(1.2), Inches(11.333), Inches(1), 'INDEXATION ET RÉFÉRENCEMENT WEB', 42, WHITE, True)
txt(s, Inches(1), Inches(2.4), Inches(11.333), Inches(0.6), 'Guide Complet et Approfondi — SEO Search Engine Optimization', 20, CYAN)
for i, tag in enumerate(['On-Page SEO', 'Techniques', 'Outils', 'Stratégie']):
    x = Inches(1 + i*2.5)
    rrect(s, x, Inches(3.4), Inches(2), Inches(0.45), RGBColor(230, 240, 255), None)
    txt(s, x, Inches(3.43), Inches(2), Inches(0.4), tag, 13, DARK_BG, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.3), Inches(11.333), Inches(0.5), 'Institut Supérieur du Numérique — SupNum, Nouakchott', 16, LIGHT_GRAY)
txt(s, Inches(1), Inches(5.8), Inches(11.333), Inches(0.5), 'Spécialité : Développement Web et Multimédia', 14, LIGHT_GRAY)
txt(s, Inches(1), Inches(6.3), Inches(11.333), Inches(0.5), 'Module : Indexation et Référencement Web', 14, LIGHT_GRAY)
txt(s, Inches(1), Inches(6.9), Inches(11.333), Inches(0.4), 'Version 1.0 — Juin 2026', 12, LIGHT_GRAY, False, PP_ALIGN.RIGHT)

# ===================== SLIDE 2: PLAN =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 2, 'Plan de la présentation — Les 8 parties du manuel', CYAN)
parts = [
    ('I', 'Fondamentaux', 'Histoire, fonctionnement, algorithmes, IA', 'Ch 1-3', CYAN),
    ('II', 'Indexation', 'Crawling, robots.txt, sitemaps, budget', 'Ch 4-7', BLUE),
    ('III', 'SEO On-Page', 'Mots-clés, intention, balises, images', 'Ch 8-12', GREEN),
    ('IV', 'SEO Technique', 'Core Web Vitals, perf, mobile, schema', 'Ch 13-18', RED),
    ('V', 'SEO Off-Page', 'Netlinking, link building, E-E-A-T', 'Ch 19-21', PURPLE),
    ('VI', 'SEO Spécialisé', 'Local, e-commerce, CMS, vidéo, tests', 'Ch 22-27', ORANGE),
    ('VII', 'Mesure', 'GSC, GA4, KPIs, reporting', 'Ch 28-30', YELLOW),
    ('VIII', 'Stratégie', 'Audit, plan action, IA, tendances', 'Ch 31-34', CYAN),
]
for i, (num, title, desc, ch, c) in enumerate(parts):
    row, col = i // 2, i % 2
    l, t = Inches(0.8 + col*6.2), Inches(1.4 + row*1.45)
    rrect(s, l, t, Inches(5.8), Inches(1.25), MID_BG, c)
    rect(s, l, t, Inches(0.7), Inches(1.25), c)
    txt(s, l+Inches(0.05), t+Inches(0.15), Inches(0.6), Inches(0.8), num, 24, WHITE, True, PP_ALIGN.CENTER)
    txt(s, l+Inches(0.85), t+Inches(0.05), Inches(3.5), Inches(0.35), title, 15, WHITE, True)
    txt(s, l+Inches(0.85), t+Inches(0.4), Inches(4.8), Inches(0.35), desc, 11, LIGHT_GRAY)
    txt(s, l+Inches(4.5), t+Inches(0.85), Inches(1.2), Inches(0.3), ch, 10, c, False, PP_ALIGN.RIGHT)

# ===================== SLIDE 3: METHODOLOGY =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 3, 'Méthodologie de rédaction', GREEN)
items = [
    ('🎯', 'Approche pédagogique', 'Structure homogène : objectifs → théorie → illustrations → exercices → synthèse\nChaque chapitre suit le même plan pour faciliter l\'apprentissage progressif'),
    ('📚', 'Sources', 'Documentation Google Search Central • Ingénieurs Google (John Mueller, Gary Illyes)\nRecherches académiques • Retours d\'expérience de campagnes SEO réelles'),
    ('🎓', 'Public & prérequis', 'Étudiants Dév. Web et Multimédia (Licence/Master) • Professionnels du numérique\nPrérequis : HTML/CSS • Notions de marketing digital (recommandé)'),
    ('📊', 'Volume horaire', '56 heures (30h cours + 26h TD/TP) • 14 semaines • 8 séquences\nÉvaluation : CC (30%) • Projet audit (30%) • Examen final (40%)'),
]
for i, (icon, title, desc) in enumerate(items):
    t = Inches(1.3 + i*1.5)
    rrect(s, Inches(0.8), t, Inches(11.733), Inches(1.35), MID_BG, GREEN, Pt(1))
    txt(s, Inches(1.1), t+Inches(0.05), Inches(0.5), Inches(0.4), icon, 22, WHITE)
    txt(s, Inches(1.7), t+Inches(0.05), Inches(3), Inches(0.35), title, 16, GREEN, True)
    txt(s, Inches(1.7), t+Inches(0.4), Inches(10.5), Inches(0.85), desc, 12, LIGHT_GRAY)

# ===================== SLIDE 4: PEDAGOGICAL =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 4, 'Structure pédagogique du manuel', GREEN)
items = [
    ('38', 'Objectifs', '3 objectifs mesurables\npar chapitre'),
    ('37', 'Synthèses', 'Points à retenir\nfin de chapitre'),
    ('4', 'Exercices', 'Mise en pratique\ndes concepts clés'),
    ('7+', 'Schémas', 'Diagrammes SVG\nillustrés'),
    ('15+', 'Outils', 'Référencés avec\nprix et usages'),
    ('4', 'Cas concrets', 'Migrations, CWV,\ncontenu, international'),
]
for i, (num, label, desc) in enumerate(items):
    row, col = i // 3, i % 3
    l, t = Inches(0.8 + col*4.1), Inches(1.5 + row*2.8)
    rrect(s, l, t, Inches(3.8), Inches(2.5), MID_BG, GREEN)
    txt(s, l, t+Inches(0.15), Inches(3.8), Inches(0.9), num, 40, GREEN, True, PP_ALIGN.CENTER)
    txt(s, l, t+Inches(1.0), Inches(3.8), Inches(0.4), label, 16, WHITE, True, PP_ALIGN.CENTER)
    txt(s, l, t+Inches(1.4), Inches(3.8), Inches(0.8), desc, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ===================== SLIDE 5: PART I-1 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 5, 'Partie I — Histoire et évolution des moteurs', CYAN)
card(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.5), 'Histoire des moteurs de recherche',
     ['1990 : Archie (1er moteur, indexe fichiers FTP)',
      '1994 : Yahoo! (annuaire humain, pas de robot)',
      '1998 : Google (PageRank, index inversé)',
      '2000 : Péripéties (bulle Internet)',
      '2004 : Mises à jour algorithmes (Florida, Boston)',
      '2011 : Panda (qualité contenu)',
      '2012 : Penguin (liens spam)',
      '2013 : Hummingbird (recherche sémantique)',
      '2015 : RankBrain (IA/ML)',
      '2019 : BERT (langage naturel)',
      '2021 : MUM (multimodal)'], CYAN)
card(s, Inches(7.2), Inches(1.5), Inches(5.3), Inches(5.5), 'Innovations clés de Google',
     ['Index inversé : mot-clé → documents',
      'PageRank : autorité via liens entrants',
      'Traitement décentralisé (GoogleFS, BigTable)',
      'Mises à jour algorithmiques régulières',
      'Parts de marché : Google 91%, Bing 3.5%',
      'Baidu (Chine) 2.5%, Yahoo 1.2%',
      'Recherche vocale : 27% des requêtes',
      'Recherche mobile > desktop depuis 2015',
      'SGE (Search Generative Experience) 2024',
      'AI Overviews déployé mondialement'], CYAN)

# ===================== SLIDE 6: PART I-2 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 6, 'Partie I — Fonctionnement technique des moteurs', BLUE)
rrect(s, Inches(0.8), Inches(1.4), Inches(11.733), Inches(1.2), MID_BG, BLUE)
txt(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8), 'Les 3 étapes fondamentales : Crawling → Indexation → Ranking', 20, WHITE, True, PP_ALIGN.CENTER)

steps = [
    ('1. Crawling (Exploration)', ['Googlebot parcourt le web via les liens',
     'Commence par une liste d\'URLs connues',
     'Extrait les liens de chaque page visitée',
     'Respecte robots.txt et meta robots',
     'Fréquence : quelques jours à semaines',
     'Le budget de crawl est limité']),
    ('2. Indexation (Stockage)', ['Analyse le contenu de chaque page',
     'Stocke dans un index gigantesque',
     'Index inversé : mot → liste de documents',
     'Prend en compte : texte, balises, attributs',
     'Ignore : mots vides (le, la, de...)',
     'Google My Business : index local']),
    ('3. Ranking (Classement)', ['200+ critères de classement',
     'PageRank : autorité transmise par liens',
     'Pertinence : correspondance requête-contenu',
     'Qualité : E-E-A-T, contenu original',
     'Performance : Core Web Vitals',
     'Personnalisation : historique, localisation']),
]
for i, (title, items) in enumerate(steps):
    l = Inches(0.8 + i*4.1)
    rrect(s, l, Inches(2.9), Inches(3.8), Inches(4.2), MID_BG, BLUE, Pt(1.5))
    bar = rect(s, l, Inches(2.9), Inches(3.8), Inches(0.45), BLUE if i == 1 else (CYAN if i == 0 else PURPLE))
    txt(s, l+Inches(0.1), Inches(2.95), Inches(3.6), Inches(0.35), title, 13, WHITE, True, PP_ALIGN.CENTER)
    multi(s, l+Inches(0.15), Inches(3.5), Inches(3.5), Inches(3.5), items, 11, LIGHT_GRAY)

# ===================== SLIDE 7: PART I-3 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 7, 'Partie I — IA et moteurs de recherche modernes', PURPLE)
algos = [
    ('RankBrain (2015)', ['1er système IA de Google',
     'Apprentissage automatique',
     'Traite les requêtes jamais vues (15%)',
     'Associe requête à concepts proches',
     'Améliore avec le temps']),
    ('BERT (2019)', ['Bidirectional Encoder Representations',
     'Comprend le contexte des mots',
     'Prépositions (à, dans, pour...) importantes',
     'Impact sur 10% des requêtes',
     'Compréhension du langage naturel']),
    ('MUM (2021)', ['Multitask Unified Model',
     'Multimodal : texte, images, vidéo',
     'Comprend 75 langues',
     'Peut générer du langage',
     'Transfère connaissance entre tâches']),
]
for i, (title, items) in enumerate(algos):
    l = Inches(0.8 + i*4.1)
    rrect(s, l, Inches(1.5), Inches(3.8), Inches(5.5), MID_BG, PURPLE)
    txt(s, l+Inches(0.2), Inches(1.6), Inches(3.4), Inches(0.35), title, 16, PURPLE, True)
    multi(s, l+Inches(0.2), Inches(2.1), Inches(3.4), Inches(3.5), items, 12, LIGHT_GRAY)

txt(s, Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.4), '💡 À retenir : Google utilise l\'IA à chaque étape — crawling, indexation ET classement', 13, YELLOW)

# ===================== SLIDE 8: PART II-1 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 8, 'Partie II — Le processus d\'indexation en profondeur', ORANGE)
rrect(s, Inches(0.8), Inches(1.4), Inches(11.733), Inches(1.5), MID_BG, ORANGE)
txt(s, Inches(1), Inches(1.45), Inches(11.3), Inches(0.4), 'Le pipeline d\'indexation', 18, ORANGE, True)
multi(s, Inches(1), Inches(1.9), Inches(11.3), Inches(0.9), [
    'Analyse → Extraction → Stockage → Mise à jour',
    'Google décompile le HTML, extrait le texte, analyse les liens, génère un index'
], 12, LIGHT_GRAY)

boxes = [
    ('Problèmes courants', ['Erreurs de crawl (DNS, timeout, 404)',
     'Soft 404 (page vide, code 200)',
     'Pages orphelines (aucun lien interne)',
     'Contenu dupliqué',
     'Pages bloquées par robots.txt',
     'Erreurs 5XX du serveur']),
    ('Solutions', ['Rapport de couverture GSC',
     'Corriger les erreurs 404 → 301',
     'Maillage interne systématique',
     'Balises canoniques',
     'Revoir robots.txt',
     'Améliorer l\'infrastructure']),
    ('Vérification', ['Google Search Console',
     'site:domain.com dans Google',
     'URL Inspection Tool',
     'Log analysis des crawls',
     'Screaming Frog crawl',
     'Rapport d\'indexation mensuel']),
]
for i, (title, items) in enumerate(boxes):
    l = Inches(0.8 + i*4.1)
    rrect(s, l, Inches(3.2), Inches(3.8), Inches(3.8), MID_BG, ORANGE, Pt(1))
    txt(s, l+Inches(0.2), Inches(3.3), Inches(3.4), Inches(0.35), title, 15, ORANGE, True)
    multi(s, l+Inches(0.2), Inches(3.75), Inches(3.4), Inches(3.2), items, 12, LIGHT_GRAY)

# ===================== SLIDE 9: PART II-2 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 9, 'Partie II — Contrôle du crawling : robots.txt et sitemaps', ORANGE)
card(s, Inches(0.8), Inches(1.4), Inches(5.8), Inches(3), 'robots.txt',
     ['Fichier texte à la racine du site',
      'Directives : User-agent, Disallow, Allow',
      'Exemple : Disallow: /admin/',
      'Ne bloque pas l\'indexation (juste le crawl)',
      'Vérification dans GSC → robots.txt Tester',
      'Attention : fichiers sensibles visibles'], ORANGE)
card(s, Inches(7.2), Inches(1.4), Inches(5.3), Inches(3), 'Sitemap XML',
     ['Liste des pages importantes du site',
      'Balise <loc>, <lastmod>, <changefreq>, <priority>',
      'Limite : 50 000 URLs par sitemap',
      'Soumission via GSC',
      'Types : news, video, image, mobile',
      'Essentiel pour les grands sites'], ORANGE)
card(s, Inches(0.8), Inches(4.6), Inches(11.733), Inches(2.5), 'Balises Meta Robots',
     ['<meta name="robots" content="noindex, nofollow">',
      'noindex : ne pas indexer la page',
      'nofollow : ne pas suivre les liens',
      ' Canonical : consolider les signaux (rel="canonical")',
      'Usage : pages duplicate, filtres, pagination',
      '⚠️ Erreur fréquente : canonical sur pagination'], ORANGE)

# ===================== SLIDE 10: PART II-3 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 10, 'Partie II — Budget de crawl et optimisation', ORANGE)
rrect(s, Inches(0.8), Inches(1.4), Inches(11.733), Inches(1.5), MID_BG, ORANGE)
txt(s, Inches(1), Inches(1.5), Inches(11.3), Inches(0.4), 'Le budget de crawl : pages explorées par période', 18, ORANGE, True)
multi(s, Inches(1), Inches(2.0), Inches(11.3), Inches(0.8), [
    'Facteurs : popularité (PageRank), fraîcheur du contenu, santé du serveur, taille du site',
    'Un budget de crawl limité signifie que Google n\'explore pas toutes vos pages'
], 12, LIGHT_GRAY)

card(s, Inches(0.8), Inches(3.2), Inches(5.8), Inches(4), 'Facteurs influençant le crawl budget',
     ['Popularité du domaine (PageRank)',
      'Nombre de pages totales',
      'Qualité du contenu (taux de rebond)',
      'Vitesse du serveur (TTFB)',
      'Fréquence de mise à jour',
      'Structure des URLs',
      'Taux d\'erreurs (4xx, 5xx)',
      'Qualité du maillage interne',
      'Présence dans le sitemap XML'], ORANGE)
card(s, Inches(7.2), Inches(3.2), Inches(5.3), Inches(4), 'Optimisation du crawl budget',
     ['Prioriser le contenu à forte valeur',
      'Supprimer les pages de faible qualité',
      'Améliorer la vitesse du serveur',
      'Optimiser le maillage interne',
      'Utiliser le sitemap XML',
      'Éviter les paramètres d\'URL',
      'Consolider le contenu dupliqué',
      'Utiliser noindex sur les pages inutiles',
      'Monitorer via le rapport Crawl Stats'], ORANGE)

# ===================== SLIDE 11: PART III-1 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 11, 'Partie III — Recherche de mots-clés (Keyword Research)', GREEN)
rrect(s, Inches(0.8), Inches(1.3), Inches(11.733), Inches(1.2), MID_BG, GREEN)
txt(s, Inches(1), Inches(1.35), Inches(11.3), Inches(0.4), 'La fondation de toute stratégie SEO', 18, GREEN, True)
multi(s, Inches(1), Inches(1.8), Inches(11.3), Inches(0.6), ['La recherche de mots-clés permet de comprendre ce que votre audience recherche et comment elle le formule'], 12, LIGHT_GRAY)

types = [
    ('Tête de réseau', 'Courts (1-2 mots)\nVolume élevé\nConcurrence forte\nEx: "SEO", "assurance"',
     'Difficulté : ★★★★★\nConversion : ★★★'),
    ('Corps', 'Moyens (2-3 mots)\nVolume moyen\nConcurrence modérée\nEx: "agence SEO Paris"',
     'Difficulté : ★★★\nConversion : ★★★★'),
    ('Longue traîne', 'Longs (3+ mots)\nVolume faible\nConcurrence faible\nEx: "agence SEO pas cher Paris"',
     'Difficulté : ★\nConversion : ★★★★★'),
]
for i, (title, desc, meta) in enumerate(types):
    l = Inches(0.8 + i*4.1)
    rrect(s, l, Inches(2.8), Inches(3.8), Inches(3.2), MID_BG, GREEN, Pt(1.5))
    txt(s, l+Inches(0.2), Inches(2.9), Inches(3.4), Inches(0.35), title, 16, GREEN, True)
    txt(s, l+Inches(0.2), Inches(3.3), Inches(3.4), Inches(1.5), desc, 12, LIGHT_GRAY)
    txt(s, l+Inches(0.2), Inches(4.9), Inches(3.4), Inches(0.8), meta, 11, YELLOW)

txt(s, Inches(0.8), Inches(6.3), Inches(11.733), Inches(0.8), 
    'Outils : Google Keyword Planner, SEMrush, Ahrefs, Ubersuggest, AnswerThePublic, AlsoAsked',
    13, LIGHT_GRAY)

# ===================== SLIDE 12: PART III-2 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 12, 'Partie III — Intention de recherche (Search Intent)', GREEN)
txt(s, Inches(1), Inches(1.3), Inches(11.333), Inches(0.5), 'Les 4 types d\'intention de recherche', 18, GREEN, True)

intents = [
    ('🔍 Informationnelle', 'L\'utilisateur cherche\nà apprendre', 'Ex: "qu\'est-ce que le SEO"\n"comment fonctionne Google"\nContenu : guides, tutoriels,\narticles explicatifs'),
    ('📍 Navigationnelle', 'L\'utilisateur cherche\nun site spécifique', 'Ex: "Facebook login"\n"YouTube"\n"Google Analytics"\nContenu : page d\'accueil'),
    ('🛍️ Commerciale', 'L\'utilisateur compare\navant d\'acheter', 'Ex: "meilleur outil SEO"\n"SEMrush vs Ahrefs"\nContenu : comparatifs,\ntests, avis'),
    ('💰 Transactionnelle', 'L\'utilisateur est\nprêt à acheter', 'Ex: "acheter nom de domaine"\n"abonnement SEMrush"\nContenu : pages produits,\nfiches prix'),
]
for i, (title, intent, example) in enumerate(intents):
    l = Inches(0.8 + i*3.1)
    rrect(s, l, Inches(1.9), Inches(2.9), Inches(3.5), MID_BG, GREEN)
    txt(s, l+Inches(0.1), Inches(2.0), Inches(2.7), Inches(0.4), title, 14, GREEN, True)
    txt(s, l+Inches(0.1), Inches(2.5), Inches(2.7), Inches(1.2), intent, 12, LIGHT_GRAY)
    txt(s, l+Inches(0.1), Inches(3.8), Inches(2.7), Inches(1.3), example, 11, YELLOW)

txt(s, Inches(0.8), Inches(5.7), Inches(11.733), Inches(0.5), 'Parcours utilisateur : Informationnel → Navigationnel → Commercial → Transactionnel', 14, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(6.2), Inches(11.733), Inches(0.5), '💡 Le contenu doit correspondre à l\'intention pour bien se classer dans les SERP', 13, YELLOW)

# ===================== SLIDE 13: PART III-3 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 13, 'Partie III — Optimisation des balises HTML', GREEN)
tags_data = [
    ('Title Tag', ['2e facteur On-Page le + important',
     '60 caractères max (affichage SERP)',
     'Unique par page',
     'Mot-clé principal au début',
     'Format : "Mot-clé - Nom du site"',
     'Éviter : keyword stuffing']),
    ('Meta Description', ['Influence le CTR (taux de clic)',
     '160 caractères max',
     'Doit donner envie de cliquer',
     'Include mot-clé + call-to-action',
     'Unique par page',
     'Peut ne pas être utilisée par Google']),
    ('Headings (H1-H6)', ['H1 = titre principal (1 par page)',
     'H2 = sections principales',
     'H3 = sous-sections',
     'Hiérarchie logique : H1 > H2 > H3',
     'Mots-clés dans les headings',
     'Accessibilité : structure claire']),
]
for i, (title, items) in enumerate(tags_data):
    l = Inches(0.8 + i*4.1)
    rrect(s, l, Inches(1.4), Inches(3.8), Inches(5.5), MID_BG, GREEN, Pt(1.5))
    txt(s, l+Inches(0.2), Inches(1.5), Inches(3.4), Inches(0.35), title, 16, GREEN, True)
    multi(s, l+Inches(0.2), Inches(2.0), Inches(3.4), Inches(3.5), items, 12, LIGHT_GRAY)

# ===================== SLIDE 14: PART III-4 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 14, 'Partie III — Images, multimédia et maillage interne', GREEN)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.8), 'Optimisation des images',
     ['Texte alternatif (alt) : essentiel accessibilité + SEO',
      'Formats modernes : WebP (-30%), AVIF (-50%)',
      'Compression : qualité 80-85% pour photos',
      'Lazy loading natif : loading="lazy"',
      'Dimensions explicites width/height (CLS)',
      'Preload pour les images LCP'], GREEN)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(2.8), 'Balise picture responsive',
     ['<picture> avec sources WebP/AVIF',
      'srcset pour différentes résolutions',
      'sizes pour adapter à la viewport',
      'Fallback JPEG pour navigateurs anciens',
      'Image CDN : Cloudinary, Imgix, imgproxy',
      'WebP supporté partout depuis 2024'], GREEN)
card(s, Inches(0.8), Inches(4.3), Inches(11.733), Inches(2.8), 'Maillage interne et Topical Clusters',
     ['Architecture de l\'information : organisation logique des contenus',
      'Topical Clusters : 1 page pilier + articles satellites interconnectés',
      'Silo structurel : catégories et sous-catégories',
      'Breadcrumbs : fil d\'Ariane pour la navigation',
      'Pages piliers : contenu long format couvrant un thème en profondeur',
      'Ancre de lien : texte descriptif (pas "cliquez ici")',
      'Profondeur max : 3 clics depuis la page d\'accueil'], GREEN)

# ===================== SLIDE 15: PART IV-1 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 15, 'Partie IV — Core Web Vitals', RED)
rrect(s, Inches(0.8), Inches(1.3), Inches(11.733), Inches(2.2), MID_BG, RED)
txt(s, Inches(1), Inches(1.4), Inches(11.3), Inches(0.4), 'Les 3 métriques de l\'expérience utilisateur (facteurs de classement depuis 2022)', 16, RED, True)
cwv = [
    ('LCP', 'Largest\nContentful\nPaint', '< 2,5s', 'Vitesse de\nchargement', GREEN),
    ('INP', 'Interaction\nto Next Paint', '< 200ms', 'Réactivité\naux clics', YELLOW),
    ('CLS', 'Cumulative\nLayout Shift', '< 0,1', 'Stabilité\nvisuelle', GOOGLE_BLUE),
]
for i, (name, full, threshold, desc, c) in enumerate(cwv):
    l = Inches(1 + i*4)
    rrect(s, l, Inches(3.8), Inches(3.5), Inches(3.2), MID_BG, c, Pt(2.5))
    txt(s, l, Inches(3.9), Inches(3.5), Inches(0.6), name, 28, c, True, PP_ALIGN.CENTER)
    txt(s, l, Inches(4.5), Inches(3.5), Inches(0.8), full, 14, WHITE, False, PP_ALIGN.CENTER)
    txt(s, l, Inches(5.4), Inches(3.5), Inches(0.5), threshold, 20, c, True, PP_ALIGN.CENTER)
    txt(s, l, Inches(5.9), Inches(3.5), Inches(0.4), desc, 11, LIGHT_GRAY, False, PP_ALIGN.CENTER)

txt(s, Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.4), 'Outils : PageSpeed Insights, Chrome UX Report, Lighthouse, CrUX Dashboard', 12, LIGHT_GRAY)

# ===================== SLIDE 16: PART IV-2 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 16, 'Partie IV — Performance web avancée', RED)
card(s, Inches(0.8), Inches(1.3), Inches(3.8), Inches(5.5), 'CDN & Caching',
     ['CDN : cache distribué géographiquement',
      'Cloudflare, Fastly, KeyCDN, BunnyCDN',
      'Cache navigateur : Cache-Control max-age',
      'Cache serveur : Varnish, NGINX, Redis',
      'Cache applicatif : APCu, Memcached',
      'Stale-while-revalidate : servir + rafraîchir',
      'TTL longs (30j+) pour assets versionnés',
      'Origin Shield : cache central'], RED)
card(s, Inches(4.9), Inches(1.3), Inches(3.8), Inches(5.5), 'HTTP & Images',
     ['HTTP/2 : multiplexage, compression headers',
      'HTTP/3 : QUIC, 0-RTT, sans TCP',
      'WebP : 30% plus léger que JPEG',
      'AVIF : 50% plus léger que JPEG',
      'Lazy loading natif : loading="lazy"',
      'Responsive images : srcset + sizes',
      'Preload : <link rel="preload"> pour LCP',
      'Image CDN : Cloudinary, Imgix'], RED)
card(s, Inches(9.0), Inches(1.3), Inches(3.8), Inches(5.5), 'JS & CSS',
     ['Code splitting : charger à la demande',
      'Tree shaking : éliminer le code mort',
      'Critical CSS : CSS nécessaire au 1er affichage',
      'Async/defer pour JS non critique',
      'Minification : HTML, CSS, JS',
      'Brotli : 20-30% mieux que Gzip',
      'TTFB < 600ms (serveur + réseau)',
      'Bundle analysis : Webpack, Vite'], RED)

# ===================== SLIDE 17: PART IV-3 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 17, 'Partie IV — Mobile-First & Données structurées', RED)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), 'Mobile-First Indexing',
     ['Google utilise la version mobile pour indexer',
      'Depuis 2019, majorité des sites en mobile-first',
      'Design responsive obligatoire',
      'Taille police minimum : 16px',
      'Espaces cliquables : 48x48px minimum',
      'Pas de contenu masqué sur mobile',
      'Test : Google Mobile-Friendly Test',
      'Performance mobile = prioritaire',
      'Core Web Vitals mobiles pénalisants'], RED)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5), 'Données structurées (Schema.org)',
     ['Vocabulaire standardisé pour décrire le contenu',
      'Formats : JSON-LD (recommandé), Microdata',
      'Types essentiels : Article, Product, FAQ, Review',
      'Rich snippets : étoiles, prix, FAQ dans SERP',
      'Test : Google Rich Results Test',
      'Erreurs fréquentes : champs manquants, mauvais type',
      'JSON-LD dans <head> ou <body>',
      'Avantage concurrentiel significatif',
      'Valider régulièrement'], RED)

# ===================== SLIDE 18: PART IV-4 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 18, 'Partie IV — JavaScript SEO', RED)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.5), 'Défis du JavaScript pour le SEO',
     ['Googlebot peut exécuter JS (Wave 2) mais avec limites',
      'Ressources JS lentes : délai d\'indexation',
      'Contenu injecté en JS non visible au crawl',
      'Erreurs JS silencieuses : contenu manquant',
      'Lazy loading excessif : pages vides au crawl',
      'Hash routing (#!) : problèmes d\'indexation'], RED)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(2.5), 'Stratégies de rendu',
     ['SSR (Server-Side Rendering) : rendu serveur, indexable',
      'SSG (Static Site Generation) : pages statiques pré-générées',
      'ISR (Incremental Static Regeneration) : SSG + MAJ dynamique',
      'Hydration : JS côté client après rendu HTML',
      'Island Architecture : îlots interactifs dans HTML statique',
      'Recommandé : SSR/SSG pour pages importantes'], RED)
card(s, Inches(0.8), Inches(4.0), Inches(11.733), Inches(3.0), 'Vérification du rendu JavaScript',
     ['Google Search Console → URL Inspection → "Afficher comme Google"',
      'Googlebot rend-il la page correctement ? Vérifier la capture',
      'Test avec "Afficher comme Google" : le contenu JS est-il visible ?',
      'Outils : Puppeteer, Playwright pour simuler Googlebot',
      'Vérifier le rendu Google via l\'URL Inspection de GSC',
       'Bon indicateur : la page dans Google est identique au rendu HTML initial',
      '⚠️ Éviter : soft 404 JS (page vide retournée avec code JS)'], RED)

# ===================== SLIDE 19: PART V-1 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 19, 'Partie V — Netlinking et Backlinks', PURPLE)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.8), 'Fondamentaux des backlinks',
     ['Les backlinks restent un facteur de classement majeur',
      'Un lien = vote de confiance d\'un site vers un autre',
      'Dofollow : transmet l\'autorité (PageRank)',
      'Nofollow : ne transmet pas (rel="nofollow")',
      'Qualité > Quantité : 1 lien .edu > 100 liens spam',
      'Profil de liens : diversité des domaines référents'], PURPLE)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(2.8), 'Qualité des backlinks',
     ['Autorité du domaine source (DR/Domain Rating)',
      'Pertinence thématique (SEO → site web, pas recettes)',
      'Position du lien (contenu principal > footer)',
      'Texte d\'ancrage (anchor text) naturel',
      'Trafic du site source (visiteurs réels)',
      'Ratio dofollow/nofollow naturel'], PURPLE)
card(s, Inches(0.8), Inches(4.3), Inches(11.733), Inches(2.8), 'Mesure et analyse des backlinks',
     ['Ahrefs : Domain Rating (DR), Referring Domains, Backlinks',
      'Majestic : Trust Flow, Citation Flow',
      'Moz : Domain Authority (DA), Page Authority (PA)',
      'Rapport mensuel : nouveaux/ perdus backlinks',
      'Désavouer les liens toxiques via Google Disavow Tool',
      'Analyse concurrentielle : qui a des backlinks que vous n\'avez pas ?'], PURPLE)

# ===================== SLIDE 20: PART V-2 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 20, 'Partie V — Link Building et E-E-A-T', PURPLE)
card(s, Inches(0.8), Inches(1.3), Inches(3.8), Inches(5.5), 'Stratégies de Link Building',
     ['Digital PR : relations presse numériques',
      'Linkable Assets : études, infographies, outils gratuits',
      'Guest Blogging : articles invités sur sites pertinents',
      'Broken Link Building : signaler des liens morts',
      'Skyscraper Technique : améliorer le contenu existant',
      'HARO (Help a Reporter Out) : citations presse',
      'Resource Pages : se faire référencer dans des listes'], PURPLE)
card(s, Inches(4.9), Inches(1.3), Inches(3.8), Inches(5.5), 'E-E-A-T (Expérience-Expertise-Autorité-Confiance)',
     ['Cadre d\'évaluation de la qualité Google',
      'Expérience : contenu basé sur l\'expérience réelle',
      'Expertise : qualifications, expertise démontrée',
      'Autorité : reconnu comme référence dans le domaine',
      'Confiance : transparence, mentions légales, contact',
      'YMYL : sites à fort impact (santé, finance, sécurité)',
      'Pages Auteur, biographies, citations'], PURPLE)
card(s, Inches(9.0), Inches(1.3), Inches(3.8), Inches(5.5), 'Améliorer son E-E-A-T',
     ['Rédiger des biographies d\'auteurs détaillées',
      'Afficher les diplômes, certifications, expériences',
      'Citer des sources officielles et académiques',
      'Maintenir les informations à jour',
      'Politique de confidentialité, CGV, mentions légales',
      'Avis clients vérifiés et réponses',
      'Contributions externes (conférences, publications)'], PURPLE)

# ===================== SLIDE 21: PART VI-1 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 21, 'Partie VI — SEO Local et Migration', ORANGE)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), 'SEO Local',
     ['Google Business Profile (GBP) : outil central',
      'Compléter TOUTES les infos (horaires, photos, services)',
      'Citations NAP : Nom, Adresse, Téléphone cohérents',
      'Annuaires : Pages Jaunes, 118000, Yelp, TripAdvisor',
      'Avis clients : répondre systématiquement',
      'Catégories GBP pertinentes',
      'Posts GBP réguliers (offres, événements)',
      'Balises Schema LocalBusiness',
      'Géociblage dans Google Search Console'], ORANGE)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5), 'Migration SEO',
     ['Types : HTTP→HTTPS, changement domaine, restructuration',
      'Checklist pré-migration : audit complet avant',
      'Mapping URL : chaque ancienne URL → nouvelle URL',
      '301 (permanent) pour les pages définitivement déplacées',
      '302 (temporaire) pour les changements temporaires',
      'Pas de boucles de redirection (A→B→A)',
      'Soumission nouveau sitemap dans GSC',
      'Suivi post-migration : J+1, J+7, J+30, J+90',
      'Plan de rollback si les performances chutent'], ORANGE)

# ===================== SLIDE 22: PART VI-2 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 22, 'Partie VI — SEO E-commerce et CMS', ORANGE)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.8), 'SEO E-commerce',
     ['Défis : contenu dupliqué (fiches produits similaires)',
      'Architecture : catégories → sous-catégories → produits',
      'Profondeur max 3 clics depuis l\'accueil',
      'Filtres et facettes : attention aux URLs infinies',
      'Avis clients : contenu unique + rich snippets',
      'Pages catégories : contenu rédactionnel + descriptions'], ORANGE)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(2.8), 'SEO CMS (WordPress, Shopify, Webflow)',
     ['WordPress : Rank Math / Yoast SEO',
      'Shopify : forces (balises auto), faiblesses (URLs, blog)',
      'Webflow : collections CMS dynamiques, sitemap auto',
      'Wix : SEO integre, limitations avancées',
      'Custom CMS : bonnes pratiques à implémenter manuellement',
      'Sécurité : mises à jour, plugins, hébergement'], ORANGE)
card(s, Inches(0.8), Inches(4.3), Inches(11.733), Inches(2.8), 'Bonnes pratiques par CMS',
     ['WordPress : permaliens personnalisés, caching (WP Rocket), images (ShortPixel), CDN (Cloudflare)',
      'Shopify : thème optimisé, apps SEO (Plug in SEO), éditer robots.txt, sitemap personnalisé',
      'Webflow : clean URLs, 301 redirects, meta balises, hosting rapide (Fastly CDN)',
      'Tous : certificat SSL, vitesse, mobile-friendly, données structurées, analytics'], ORANGE)

# ===================== SLIDE 23: PART VI-3 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 23, 'Partie VI — SEO Vidéo, Vocal et A/B Testing', ORANGE)
card(s, Inches(0.8), Inches(1.3), Inches(3.8), Inches(5.5), 'SEO Vidéo',
     ['YouTube = 2e moteur de recherche mondial',
      'Titre : mot-clé principal au début',
      'Description : 200+ mots, liens, timestamps',
      'Tags : 5-10 tags pertinents',
      'Miniature : personnalisée, texte lisible',
      'Sous-titres : transcription complète',
      'Données structurées VideoObject',
      'Playlists : regroupement thématique',
      'Engagement : likes, commentaires, partages'], ORANGE)
card(s, Inches(4.9), Inches(1.3), Inches(3.8), Inches(5.5), 'SEO Vocal (Voice Search)',
     ['27% des recherches sont vocales',
      'Featured snippets = source principale',
      'Réponses concises : 30-40 mots',
      'Format questions-réponses',
      'Langage naturel et conversationnel',
      'PageSpeed : sites rapides favorisés',
      'HTTPS obligatoire',
      'Données structurées FAQ/HowTo',
      '"Position zéro" : snippet + vocal'], ORANGE)
card(s, Inches(9.0), Inches(1.3), Inches(3.8), Inches(5.5), 'A/B Testing SEO',
     ['SEO Split Testing : tester sans pénaliser',
      'Durée minimum : 2-4 semaines (cycles Google)',
      'Éléments à tester : title tags, meta desc, H1',
      'Contenu : longueur, structure, mots-clés',
      'Données structurées : FAQ, Product, Review',
      'Signification statistique : p < 0.05',
      'Outils : Google Optimize, SEOTesting.com',
      'Une seule variable à la fois',
      'Splitting par URLs (pas aléatoire)'], ORANGE)

# ===================== SLIDE 24: PART VII-1 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 24, 'Partie VII — Google Search Console (GSC)', YELLOW)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.5), 'Rapports essentiels GSC',
     ['Rapport de couverture : pages indexées, erreurs, exclues',
      'Rapport de performance : clics, impressions, CTR, position',
      'Inspection d\'URL : test d\'indexation par URL',
      'Sitemaps : soumission et statut',
      'robots.txt : testeur de directives',
      'Pages mobiles : problèmes d\'ergonomie mobile'], YELLOW)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(2.5), 'Fonctionnalités avancées',
     ['Rapport Core Web Vitals : LCP, INP, CLS par URL',
      'Rapport de liage : pages les plus liées en interne',
      'Rapport données structurées : erreurs schema',
      'Rapport manuel actions : pénalités manuelles',
      'Messages Google : notifications importantes',
      'Paramètres d\'URL : gestion des paramètres'], YELLOW)
card(s, Inches(0.8), Inches(4.0), Inches(11.733), Inches(3.0), 'KPIs essentiels à suivre dans GSC',
     ['Clics : nombre de visites depuis la recherche Google',
      'Impressions : nombre d\'affichages dans les résultats',
      'CTR : clics / impressions (moyenne : 3-5%)',
      'Position moyenne : 1-10 (bon), 10-30 (à améliorer)',
      'Pages indexées : surveillance hebdomadaire',
      'Erreurs 404 : repérer et corriger rapidement',
      'Growth : tendances hebdomadaires et mensuelles'], YELLOW)

# ===================== SLIDE 25: PART VII-2 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 25, 'Partie VII — Google Analytics 4 et KPIs', YELLOW)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), 'Google Analytics 4 (GA4)',
     ['Modèle événementiel (pas de sessions)',
      'Engagement : engaged sessions, engagement rate',
      'Trafic organique : Users, New Users, Sessions',
      'Pages et écrans : pages les plus visitées',
      'Acquisition : trafic par canal (organic, direct, social)',
      'Entonnoirs de conversion : parcours utilisateur',
      'Événements : scroll, click, video, conversion',
      'BigQuery : export des données brutes',
      'GA4 vs Universal : migration obligatoire faite'], YELLOW)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5), 'KPIs SEO et Reporting',
     ['Trafic organique : volume et croissance mensuelle',
      'Positions : suivi des mots-clés cibles',
      'Taux de conversion SEO : objectifs atteints',
      'ROI SEO : valeur du trafic organique',
      'Domain Rating / Authority : évolution mensuelle',
      'Core Web Vitals : LCP, INP, CLS (bon/pass/échec)',
      'Backlinks : nouveaux domaines référents',
      'Pages indexées : % pages indexées / total crawlées',
      'Tableau de bord : Data Studio, Looker, Sheets'], YELLOW)

# ===================== SLIDE 26: PART VIII-1 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 26, 'Partie VIII — Audit SEO et Plan d\'action', CYAN)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), 'Audit SEO Complet (3 phases)',
     ['Phase 1 — Audit technique',
      '  • Crawl avec Screaming Frog (500 URLs gratuit)',
      '  • Robots.txt, sitemap.xml, HTTPS, redirections',
      '  • Core Web Vitals, mobile-friendly',
      '  • Données structurées, balises HTML',
      '',
      'Phase 2 — Audit On-Page',
      '  • Title tags, meta descriptions, headings',
      '  • Mots-clés et pertinence du contenu',
      '  • Qualité rédactionnelle, unicité',
      '',
      'Phase 3 — Audit Off-Page',
      '  • Profil de backlinks (Ahrefs, Majestic)',
      '  • E-E-A-T, mentions de marque',
      '  • Analyse concurrentielle'], CYAN)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(5.5), 'Plan d\'Action Stratégique',
     ['Priorisation par impact et effort (matrix ICE)',
      '',
      'Quick Wins (1-2 sem.) :',
      '  • Erreurs 404, meta descriptions',
      '  • Optimisation images, lazy loading',
      '',
      'Moyen terme (1-2 mois) :',
      '  • Contenu piliers, topical clusters',
      '  • Stratégie de link building',
      '',
      'Long terme (3-6 mois) :',
      '  • Migration/restructuration site',
      '  • Stratégie de contenu complète',
      '  • Amélioration E-E-A-T'],
     CYAN)

# ===================== SLIDE 27: PART VIII-2 =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 27, 'Partie VIII — SEO, IA et Tendances 2026-2030', CYAN)
card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(2.5), 'SEO et Intelligence Artificielle',
     ['Google SGE (Search Generative Experience) : réponses générées par IA',
      'AI Overviews : résumés IA en haut des SERP',
      'Impact : baisse du trafic vers les sites (zero-click)',
      'IA pour le SEO : génération de contenu, briefs',
      'Analyse sémantique : compréhension du contexte',
      'Personalisation : résultats adaptés à chaque utilisateur'], CYAN)
card(s, Inches(7.2), Inches(1.3), Inches(5.3), Inches(2.5), 'Tendances émergentes',
     ['Zero-Click Searches : réponses directes sans clic',
      'Search Everywhere : SEO sur toutes les plateformes',
      'Recherche visuelle : Google Lens, images',
      'Recherche vocale : 50% des recherches en 2026',
      'Web3 et SEO décentralisé (nouveaux moteurs)',
      'EEAT renforcé : importance croissante'], CYAN)
card(s, Inches(0.8), Inches(4.0), Inches(11.733), Inches(3.0), 'Prévisions 2026-2030',
     ['L\'IA générative va transformer les SERP : moins de clics, plus de réponses directes',
      'Le SEO devient "Search Everywhere Optimization" (Amazon, YouTube, TikTok, ChatGPT)',
      'L\'expérience utilisateur (UX) devient le facteur différenciateur principal',
      'Les données structurées et le contenu sémantique seront indispensables',
      'La qualité et l\'originalité du contenu primeront sur la quantité',
      'Recommandation : diversifier les sources de trafic, investir dans la marque'], CYAN)

# ===================== SLIDE 28: CASE STUDIES =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 28, 'Études de cas — Résultats concrets', GREEN)
cases = [
    ('Migration site', 'HTTP→HTTPS + changement domaine', 'Trafic : -30% (J1)\nRécupération : J+45\n→ Trafic x2 à M+6', GREEN),
    ('E-commerce CWV', 'Optimisation LCP, CLS, INP', 'LCP : 5,8s → 1,9s\nCLS : 0,35 → 0,05\nINP : 450ms → 120ms\n+15% trafic', BLUE),
    ('Stratégie contenu', 'Topical authority blog', '50 → 120 articles\nTrafic x5 en 12 mois\n+350% backlinks\nMeilleures positions', PURPLE),
    ('SEO International', 'Expansion 5 marchés', '+300% trafic FR\nStructure sous-rép.\nAdaptation culturelle\n30% contenu original', ORANGE),
]
for i, (title, subtitle, results, color) in enumerate(cases):
    row, col = i // 2, i % 2
    l, t = Inches(0.8 + col*6.2), Inches(1.4 + row*2.9)
    rrect(s, l, t, Inches(5.8), Inches(2.7), MID_BG, color, Pt(2))
    bar = rect(s, l, t, Inches(5.8), Inches(0.5), color)
    txt(s, l+Inches(0.2), t+Inches(0.03), Inches(3.5), Inches(0.4), title, 17, DARK_BG, True)
    txt(s, l+Inches(4), t+Inches(0.03), Inches(1.6), Inches(0.4), subtitle, 10, DARK_BG, False, PP_ALIGN.RIGHT)
    txt(s, l+Inches(0.3), t+Inches(0.7), Inches(5.3), Inches(1.8), results, 13, WHITE)

# ===================== SLIDE 29: APPENDICES =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
slide_header(s, 29, 'Annexes et ressources complémentaires', PURPLE)
items = [
    ('📋 Annexe A — Checklist technique SEO',
     ['7 catégories : Crawling, Architecture, Performance, Mobile, Schema, Sécurité, International',
      '30+ points de vérification avec seuils précis (LCP < 2.5s, CLS < 0.1, etc.)']),
    ('🛠️ Annexe B — Guide des outils SEO',
     ['15 outils avec prix : Google Search Console (gratuit), Screaming Frog (gratuit/£149), Ahrefs ($99/mois)',
      'Métriques clés : trafic organique, positions, CTR, taux d\'indexation, DR, CWV, backlinks']),
    ('🌐 Annexe C — Ressources',
     ['Veille : Google Search Status, SEOFrog Roadmap, Twitter @googlesearchc',
      'Communautés : WebmasterWorld, Reddit r/SEO, Google Search Central Help',
      'Certifications : Google SEO Fundamentals, Ahrefs Academy, SEMrush Academy']),
]
for i, (title, desc) in enumerate(items):
    t = Inches(1.5 + i*1.9)
    rrect(s, Inches(0.8), t, Inches(11.733), Inches(1.6), MID_BG, PURPLE)
    txt(s, Inches(1.1), t+Inches(0.1), Inches(11.2), Inches(0.35), title, 15, PURPLE, True)
    multi(s, Inches(1.1), t+Inches(0.5), Inches(11.2), Inches(1.0), desc, 12, LIGHT_GRAY)

# ===================== SLIDE 30: THANK YOU =====================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_BG)
rect(s, Inches(0), Inches(0), Inches(13.333), Inches(4.5), MID_BG)
accent(s, Inches(2), Inches(4.5), Inches(9.333), CYAN)
dot_pattern(s)
txt(s, Inches(1), Inches(1.5), Inches(11.333), Inches(1.2), 'Merci de votre attention', 44, WHITE, True, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3), Inches(11.333), Inches(0.8), 'Guide complet disponible sur GitHub', 22, CYAN, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.2), Inches(11.333), Inches(0.5), 'Institut Supérieur du Numérique — SupNum, Nouakchott', 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.7), Inches(11.333), Inches(0.5), 'Développement Web et Multimédia — Indexation et Référencement Web', 14, LIGHT_GRAY, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5), 'github.com/mohameden19961/rapport-seo', 16, BLUE, False, PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(6.8), Inches(11.333), Inches(0.4), 'Version 1.0 — Juin 2026', 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Save
output_path = '/home/abdy/rapport_seo/presentation_seo.pptx'
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Size: {os.path.getsize(output_path) / 1024:.1f} KB")
print(f"Slides: {len(prs.slides)}")
